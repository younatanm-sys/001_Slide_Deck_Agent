"""Core skill for generating PowerPoint slides."""

from typing import Optional, List
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from ..models import SlideContent, SlideType, PresentationRequest, GenerationResult


class SlideGeneratorSkill:
    """Skill for generating PowerPoint presentations."""

    def __init__(self):
        """Initialize the slide generator skill."""
        self.prs: Optional[Presentation] = None

    def create_presentation(self, request: PresentationRequest) -> GenerationResult:
        """
        Create a complete presentation from a request.

        Args:
            request: PresentationRequest with all presentation details

        Returns:
            GenerationResult with success status and metadata
        """
        try:
            self.prs = Presentation()
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)

            # Set metadata
            if request.author:
                self.prs.core_properties.author = request.author
            if request.company:
                self.prs.core_properties.category = request.company
            self.prs.core_properties.title = request.topic

            # Generate slides
            for slide_content in request.slides:
                self._add_slide(slide_content, request)

            # Save presentation
            output_path = Path(request.output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(str(output_path))

            return GenerationResult(
                success=True,
                output_path=str(output_path),
                slide_count=len(request.slides),
                metadata={
                    "topic": request.topic,
                    "template": request.template,
                    "author": request.author,
                },
            )

        except Exception as e:
            return GenerationResult(
                success=False, error=str(e), metadata={"request": request.model_dump()}
            )

    def _add_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a slide based on its type."""
        if content.slide_type == SlideType.TITLE:
            self._add_title_slide(content, request)
        elif content.slide_type == SlideType.TITLE_CONTENT:
            self._add_title_content_slide(content, request)
        elif content.slide_type == SlideType.SECTION_HEADER:
            self._add_section_header_slide(content, request)
        elif content.slide_type == SlideType.TWO_COLUMN:
            self._add_two_column_slide(content, request)
        elif content.slide_type == SlideType.BULLET_POINTS:
            self._add_bullet_points_slide(content, request)
        elif content.slide_type == SlideType.QUOTE:
            self._add_quote_slide(content, request)
        elif content.slide_type == SlideType.THANK_YOU:
            self._add_thank_you_slide(content, request)
        else:
            self._add_blank_slide()

    def _add_title_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(request.primary_color)

        # Title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.text = content.title or "Untitled Presentation"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Subtitle
        if content.subtitle:
            subtitle_top = Inches(4.2)
            subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1))
            subtitle_frame = subtitle_box.text_frame

            p = subtitle_frame.paragraphs[0]
            p.text = content.subtitle
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(230, 230, 230)

        # Add speaker notes if provided
        if content.notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = content.notes

    def _add_title_content_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a title and content slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = content.title or ""
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(request.primary_color)

        # Add line under title
        line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.4), Inches(9), Inches(0.02)  # Line shape
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(request.secondary_color)
        line.line.fill.background()

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(8.6), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = MSO_ANCHOR.TOP

        p = content_frame.paragraphs[0]
        p.text = content.content or ""
        p.font.size = Pt(18)
        p.font.color.rgb = self._hex_to_rgb(request.text_color)
        p.space_after = Pt(12)

        if content.notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = content.notes

    def _add_section_header_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a section header slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(request.secondary_color)

        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame

        p = title_frame.paragraphs[0]
        p.text = content.title or ""
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    def _add_two_column_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a two-column slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = content.title or ""
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(request.primary_color)

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(4.1), Inches(5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        p = left_frame.paragraphs[0]
        p.text = content.left_content or ""
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(request.text_color)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.1), Inches(5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        p = right_frame.paragraphs[0]
        p.text = content.right_content or ""
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(request.text_color)

    def _add_bullet_points_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a bullet points slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = content.title or ""
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(request.primary_color)

        # Bullet points
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(8.6), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        if content.bullet_points:
            for i, bullet in enumerate(content.bullet_points):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()

                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = self._hex_to_rgb(request.text_color)
                p.space_after = Pt(12)

    def _add_quote_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a quote slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Quote text
        quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(2))
        quote_frame = quote_box.text_frame
        quote_frame.word_wrap = True

        p = quote_frame.paragraphs[0]
        p.text = f'"{content.quote_text}"' if content.quote_text else ""
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = self._hex_to_rgb(request.text_color)

        # Author
        if content.quote_author:
            author_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(0.5))
            author_frame = author_box.text_frame
            p = author_frame.paragraphs[0]
            p.text = f"— {content.quote_author}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.color.rgb = self._hex_to_rgb(request.secondary_color)

    def _add_thank_you_slide(self, content: SlideContent, request: PresentationRequest):
        """Add a thank you slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(request.primary_color)

        # Thank you text
        text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        text_frame = text_box.text_frame

        p = text_frame.paragraphs[0]
        p.text = content.title or "Thank You"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Subtitle (e.g., contact info)
        if content.subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = content.subtitle
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(230, 230, 230)

    def _add_blank_slide(self):
        """Add a blank slide."""
        self.prs.slides.add_slide(self.prs.slide_layouts[6])

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        )
