"""
Main Slide Generator - v2.0
============================

Professional slide generator implementing design specifications v2.0 with:
- PRINCIPLE 1: Zoned Integrity (Title/Content/Footer zones immutable)
- PRINCIPLE 2: Content-Driven Layouts (Split-Screen 60/40, Full-Width, Hero-Visual)
- PRINCIPLE 3: Atomic Text Elements (each bullet = separate paragraph)
- PRINCIPLE 4: Hierarchical Titling (Slide title T1 + Chart title T2)
- PRINCIPLE 5: Consistent Visual Anchoring (vertical center titles, top-align content)
- Dynamic layout engine (chart+insight, single component, bullets)
- Chart generation (column, waterfall, matrix)
- Chart annotations (CAGR arrows, difference lines, callouts)
- Configurable highlight logic
- Proper 16:9 aspect ratio (1920x1080px at 144 DPI)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, BubbleChartData
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pathlib import Path
from typing import Dict, List, Any, Optional
import math

from ..models import PresentationRequest, GenerationResult
from ..templates.main_template_config import (
    MAIN_CONFIG,
    get_safe_zone_bounds,
    get_region,
    get_typography,
    validate_bounds,
    px_to_inches
)


class MainSlideGeneratorSkill:
    """
    Main slide generator implementing v2.0 design specifications.

    Features:
    - Dynamic layout engine that analyzes content and chooses appropriate layout
    - Chart generation with professional styling (column charts, bar charts)
    - Proper text vertical centering in all regions
    - Professional bullet formatting with hanging indents
    - Enforces Safe Zone boundaries (96,54) to (1824,1026)
    - Applies typography hierarchy (T1-T5) with overflow logic
    - Uses proper 16:9 aspect ratio (13.33" x 7.5")
    """

    def __init__(self):
        """Initialize Main slide generator with v2.0 config."""
        self.config = MAIN_CONFIG
        self.grid = self.config['grid']
        self.typography = self.config['typography']
        self.colors = self.config['colors']

    def create_presentation(self, request: PresentationRequest) -> GenerationResult:
        """
        Create a professionally-styled presentation following v2.0 specifications.

        Args:
            request: Presentation request with topic and slides

        Returns:
            GenerationResult with success status and metadata
        """
        try:
            # Create presentation with proper 16:9 dimensions
            prs = Presentation()
            prs.slide_width = self.grid['canvas']['width']  # 13.33"
            prs.slide_height = self.grid['canvas']['height']  # 7.5"

            # Add title slide
            self._add_title_slide(prs, request)

            # Add content slides with dynamic layout
            for slide_content in request.slides:
                self._add_content_slide(prs, slide_content)

            # Add closing slide
            self._add_closing_slide(prs, request)

            # Save presentation
            output_path = Path(request.output_path)
            prs.save(str(output_path))

            return GenerationResult(
                success=True,
                output_path=str(output_path),
                slide_count=len(prs.slides),
                metadata={
                    "template": "Main v2.0",
                    "aspect_ratio": "16:9",
                    "dimensions": "1920x1080px",
                    "safe_zone": "1728x972px",
                    "grid_system": "enforced"
                }
            )

        except Exception as e:
            return GenerationResult(
                success=False,
                error=str(e),
                output_path=request.output_path
            )

    def _add_title_slide(self, prs: Presentation, request: PresentationRequest):
        """Add styled title slide with full-bleed green background."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Set background to Primary Green
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.colors['primary_green_rgb'])

        # Get typography for title slide
        title_spec = self.config['typography_special']['title_slide_main']
        subtitle_spec = self.config['typography_special']['title_slide_subtitle']

        # Add main title (centered both axes)
        title_box = slide.shapes.add_textbox(
            Inches(0),
            Inches(2.5),
            self.grid['canvas']['width'],
            Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = request.topic
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Format title
        p = title_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.name = 'Arial'
        p.font.size = title_spec['font_size']
        p.font.bold = True
        p.font.color.rgb = RGBColor(*title_spec['color_rgb'])

        # Add subtitle if company provided
        if request.company:
            subtitle_y = Inches(2.5) + Inches(2) + subtitle_spec['spacing_from_title']
            subtitle_box = slide.shapes.add_textbox(
                Inches(0),
                subtitle_y,
                self.grid['canvas']['width'],
                Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = request.company
            subtitle_frame.word_wrap = True
            subtitle_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            # Format subtitle
            p = subtitle_frame.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.font.name = 'Arial'
            p.font.size = subtitle_spec['font_size']
            p.font.bold = False
            p.font.color.rgb = RGBColor(*subtitle_spec['color_rgb'])

    def _add_content_slide(self, prs: Presentation, slide_content):
        """
        Add styled content slide with dynamic layout engine.

        This function analyzes the slide content and selects the appropriate layout:
        - If chart data exists: Use chart_plus_insight_50_50 layout
        - Otherwise: Use bullets-only layout
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Set white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.colors['white_rgb'])

        # Get layout regions
        title_region = get_region('title')
        content_region = get_region('content')

        # Add title with vertical centering
        self._add_slide_title(slide, slide_content.title, title_region)

        # DYNAMIC LAYOUT ENGINE: Analyze content and choose layout
        has_chart_data = self._has_chart_data(slide_content)

        if has_chart_data:
            # Use chart+insight (60/40) layout
            self._add_chart_insight_layout(slide, slide_content)
        else:
            # Use bullets-only layout
            self._add_bullets_layout(slide, slide_content, content_region)

        # Add footer if source available
        self._add_footer(slide, slide_content)

    def _add_slide_title(self, slide, title_text: str, title_region: dict):
        """
        Add slide title with proper vertical centering.

        PRINCIPLE 1 - Zoned Integrity: Title is confined to Title zone, cannot
        violate boundaries or overlap with Content zone.

        PRINCIPLE 4 - Hierarchical Titling: This is the primary T1 title that
        communicates the key insight or "so what" of the slide.

        PRINCIPLE 5 - Consistent Visual Anchoring: Title text is vertically
        centered within the Title zone for polished, balanced appearance.

        Args:
            slide: PowerPoint slide object
            title_text: Title text
            title_region: Title region specification
        """
        title_spec = get_typography('T1')

        # PRINCIPLE 1: Title constrained to Title zone bounds
        title_box = slide.shapes.add_textbox(
            title_region['bounds']['x1'],
            title_region['bounds']['y1'],
            px_to_inches(title_region['bounds']['x2_px'] - title_region['bounds']['x1_px']),
            title_region['height']
        )

        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # PRINCIPLE 5: Vertical center
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.margin_bottom = Inches(0)

        # Format title (T1: 28pt Bold Primary Green)
        p = title_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = 'Arial'
        p.font.size = title_spec['font_size']
        p.font.bold = True
        p.font.color.rgb = RGBColor(*title_spec['color_rgb'])

    def _has_chart_data(self, slide_content) -> bool:
        """
        Analyze slide content to determine if it contains chart data.

        Returns:
            True if chart data exists, False otherwise
        """
        # Check for chart_data attribute
        if hasattr(slide_content, 'chart_data') and slide_content.chart_data:
            return True

        # Check for data attribute (alternative naming)
        if hasattr(slide_content, 'data') and slide_content.data:
            return True

        return False

    def _add_chart_insight_layout(self, slide, slide_content):
        """
        Add chart+insight layout (60% chart, 40% bullets).

        PRINCIPLE 1 - Zoned Integrity: Chart and text confined to Content zone,
        never violating Title or Footer zones.

        PRINCIPLE 2 - Content-Driven Layouts: Split-Screen Layout (60/40) used
        when content includes both visualization and explanatory text. Chart and
        text occupy separate, non-overlapping columns within Content zone.

        PRINCIPLE 4 - Hierarchical Titling: Chart must have descriptive T2 title
        explaining "what it is" below the main slide T1 title "so what".

        PRINCIPLE 5 - Consistent Visual Anchoring: Chart and text top-aligned
        within Content zone for consistent horizontal sightline.

        IMPROVEMENT 1 - Gutter-First Approach: Mandatory 40px gutter between chart
        and text (at 96 DPI final render). Config uses 60px at 144 DPI reference.
        Text box MUST start at chart_area.x2 + 60px. Non-negotiable.

        Args:
            slide: PowerPoint slide object
            slide_content: Slide content with chart data and bullets
        """
        # Get layout specification (PRINCIPLE 2: Split-Screen 60/40)
        layout = self.config['layout_distribution']['chart_plus_insight_50_50']
        chart_area = layout['chart_area']
        text_area = layout['text_area']
        gutter_px = layout['gutter_px']  # IMPROVEMENT 1: Explicit gutter (60px at 144 DPI = 40px at 96 DPI)

        # IMPROVEMENT 1: Verify gutter enforcement
        # Text area MUST start at chart_area.x2 + gutter (60px at 144 DPI config)
        assert text_area['x1_px'] == chart_area['x2_px'] + gutter_px, \
            f"Gutter violation: text starts at {text_area['x1_px']}, should be {chart_area['x2_px'] + gutter_px}"

        # Get chart data
        chart_data = slide_content.chart_data if hasattr(slide_content, 'chart_data') else slide_content.data

        # CHART TITLE SKILL: Mandatory T2 title for all charts (consulting standard)
        # This is part of the chart element for vertical centering calculations
        chart_title_height = px_to_inches(40)  # T2 20pt + spacing
        chart_has_title = isinstance(chart_data, dict) and 'title' in chart_data and chart_data['title']

        # TRUE VERTICAL CENTERING: Calculate total element height (chart + title)
        # Formula: Y_position = Content_Zone_Y + (Content_Zone_Height - Total_Height) / 2

        # Work with EMU values for precision, convert back to Inches
        content_zone_height_emu = int(chart_area['height'])  # EMU
        chart_title_height_emu = int(chart_title_height)  # EMU

        # Estimate chart body height (will be refined after creation)
        chart_body_height_emu = int(content_zone_height_emu * 0.7)  # 70% of available

        if chart_has_title:
            total_element_height_emu = chart_title_height_emu + chart_body_height_emu
        else:
            total_element_height_emu = chart_body_height_emu

        # Calculate centered position within content zone
        vertical_whitespace_emu = content_zone_height_emu - total_element_height_emu
        vertical_offset_emu = int(vertical_whitespace_emu / 2)

        # Position chart element (title + chart) at centered Y
        chart_element_top_emu = int(chart_area['y1']) + vertical_offset_emu
        chart_element_top = Inches(chart_element_top_emu / 914400)

        # Convert heights to Inches for shape creation
        chart_body_height_estimate = Inches(chart_body_height_emu / 914400)

        # Add chart title if present
        if chart_has_title:
            self._add_chart_title(
                slide,
                chart_data['title'],
                chart_area['x1'],
                chart_element_top,
                chart_area['width'],
                chart_title_height
            )
            chart_top = chart_element_top + chart_title_height
            chart_height = chart_body_height_estimate
        else:
            chart_top = chart_element_top
            chart_height = chart_body_height_estimate

        # LEFT COLUMN: Create chart based on type (PRINCIPLE 2: separate column)
        # FILL THE SPACE: Scale chart to 90% of container width for commanding presence
        chart_scale = 0.90  # 90% fill to be dominant visual element

        # Get width as EMU value, scale it, then convert back to Inches
        chart_width_emu = int(chart_area['width'])  # EMU value
        chart_scaled_width_emu = int(chart_width_emu * chart_scale)
        chart_x_offset_emu = int((chart_width_emu - chart_scaled_width_emu) / 2)

        # Convert to Inches objects for add_chart
        chart_scaled_width = Inches(chart_scaled_width_emu / 914400)
        chart_left = Inches((int(chart_area['x1']) + chart_x_offset_emu) / 914400)

        chart_type = chart_data.get('type', 'column') if isinstance(chart_data, dict) else 'column'

        if chart_type == 'waterfall':
            chart_shape = self._create_waterfall_chart(
                slide, chart_data, chart_left, chart_top, chart_scaled_width, chart_height
            )
        elif chart_type == 'matrix':
            chart_shape = self._create_matrix_chart(
                slide, chart_data, chart_left, chart_top, chart_scaled_width, chart_height
            )
        else:  # Default to column chart
            chart_shape = self._create_column_chart(
                slide, chart_data, chart_left, chart_top, chart_scaled_width, chart_height
            )

        # Add chart annotations if provided
        if isinstance(chart_data, dict) and 'annotations' in chart_data:
            # GLOBAL FIX: Pass actual chart bounds (after centering/scaling), not layout container
            actual_chart_bounds = {
                'x1': chart_left,
                'y1': chart_top,
                'width': chart_scaled_width,
                'height': chart_height
            }
            self._add_chart_annotations(slide, chart_shape, chart_data['annotations'], actual_chart_bounds)

        # FINAL PROFESSIONAL STANDARD 4: Chart Source Annotation Skill
        # If source is provided, place it in bottom-left corner of chart area (not slide footer)
        if isinstance(chart_data, dict) and 'source' in chart_data and chart_data['source']:
            self._add_chart_source(slide, chart_data['source'], chart_area)

        # IMPROVEMENT 5: Add key takeaway box if provided
        if hasattr(slide_content, 'key_takeaway') and slide_content.key_takeaway:
            # Position in lower right of chart area by default
            takeaway_x = chart_area['x1'] + chart_area['width'] * 0.05  # 5% from left edge
            takeaway_y = chart_area['y1'] + chart_area['height'] * 0.75  # 75% down
            self._add_key_takeaway_box(slide, slide_content.key_takeaway, takeaway_x, takeaway_y)

        # RIGHT COLUMN: Add bullet points (insights) - PRINCIPLE 2: separate column
        # PRINCIPLE 3: Each bullet as atomic element
        # MASTER VERTICAL CENTERING: Align text relative to chart, not content zone
        bullets = slide_content.bullet_points or []
        if bullets:
            # Calculate chart's vertical position and height
            chart_vertical_center = chart_top + (chart_height / 2)

            # Text block should be centered relative to chart's vertical center
            # Use full text area height, but anchor to MIDDLE so it aligns with chart
            text_block_top = text_area['y1']
            text_block_height = text_area['height']

            self._add_bullet_textbox(
                slide,
                bullets,
                text_area['x1'],
                text_block_top,
                text_area['width'],
                text_block_height,  # Use full height for flexibility
                vertical_anchor='MIDDLE',  # Center content within the text box
                chart_center_y=chart_vertical_center  # Pass chart center for alignment
            )

    def _add_bullets_layout(self, slide, slide_content, content_region: dict):
        """
        Add bullets-only layout occupying full content area.

        PRINCIPLE 1 - Zoned Integrity: Bullets confined to Content zone, never
        violating Title or Footer zones.

        PRINCIPLE 2 - Content-Driven Layouts: Full-Width Text Layout used when
        content is purely textual without charts.

        PRINCIPLE 3 - Atomic Text Elements: Each bullet rendered as separate
        paragraph with complete formatting.

        PRINCIPLE 5 - Consistent Visual Anchoring: Bullets top-aligned within
        Content zone for consistent horizontal sightline.

        Args:
            slide: PowerPoint slide object
            slide_content: Slide content with bullets
            content_region: Content region specification
        """
        bullets = slide_content.bullet_points or []
        if not bullets and slide_content.content:
            bullets = [slide_content.content]

        # PRINCIPLE 1, 2: Use full Content zone for text-only layout
        if bullets:
            self._add_bullet_textbox(
                slide,
                bullets,
                content_region['bounds']['x1'],
                content_region['bounds']['y1'],
                px_to_inches(content_region['bounds']['x2_px'] - content_region['bounds']['x1_px']),
                content_region['height']
            )

    def _add_bullet_textbox(self, slide, bullets: List[str], left, top, width, height, vertical_anchor='TOP', chart_center_y=None):
        """
        Add text box with properly formatted bullets.

        PRINCIPLE 3 - Atomic Text Elements: Each bullet point is generated as its own
        separate paragraph with complete formatting rules (bullet character, hanging
        indent, line spacing) to guarantee scannability and prevent dense paragraphs.

        PRINCIPLE 5 - Consistent Visual Anchoring: Content is top-aligned within the
        Content zone to maintain consistent horizontal sightline.

        MASTER VERTICAL CENTERING - Text blocks are vertically centered relative to
        the chart element, not the content zone. This ensures balanced composition
        where elements align to each other's visual weight.

        Args:
            slide: PowerPoint slide object
            bullets: List of bullet points
            left, top, width, height: Position and size in inches
            vertical_anchor: 'TOP' (default) or 'MIDDLE' for centering
            chart_center_y: Optional chart vertical center for relative alignment
        """
        body_spec = get_typography('T3')

        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # IMPROVEMENT 3: Support both TOP and MIDDLE anchoring
        if vertical_anchor == 'MIDDLE':
            content_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        else:
            content_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # PRINCIPLE 5: Default

        content_frame.margin_left = Inches(0.2)
        content_frame.margin_right = Inches(0.2)
        content_frame.margin_top = Inches(0.3)
        content_frame.margin_bottom = Inches(0.1)

        # PRINCIPLE 3: Each bullet as atomic element
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()

            # Set bullet text
            p.text = bullet

            # FILL THE SPACE: Full justification for clean block shape
            p.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY

            # Apply professional bullet formatting (hanging indent)
            self._apply_bullet_formatting(p, level=0)

            # Apply typography (T3: 18pt Body Text - increased for legibility)
            p.font.name = 'Arial'
            p.font.size = body_spec['font_size']
            p.font.bold = False
            p.font.color.rgb = RGBColor(*body_spec['color_rgb'])

            # FINAL PROFESSIONAL STANDARD 2: Mandatory 12pt line spacing after bullets
            # This creates clean, consistent rhythm down the page and uses vertical space effectively
            p.space_after = Pt(12)

    def _apply_bullet_formatting(self, paragraph, level: int = 0):
        """
        Apply professional bullet formatting with THREE SPECIFIC RULES.

        IMPROVEMENT 2 - Professional Bullet:
        Rule A: Bullet character is solid circle at 80-90% of text size,
                colored Primary Green (#147B58)
        Rule B: Hanging indent - wrapped text aligns with first line text,
                not the bullet character
        Rule C: Text gutter - mandatory 10-12px gap between bullet and text

        Args:
            paragraph: Text paragraph to format
            level: Indentation level (0, 1, 2)
        """
        # Get bulleted list config
        bullet_config = self.config['text_components']['bulleted_list']

        # Set bullet level
        paragraph.level = level

        # RULE A: Bullet character (solid circle) with correct size and color
        # python-pptx doesn't have a direct bullet API, so we access the XML element
        from lxml import etree
        pPr = paragraph._element.get_or_add_pPr()

        # Remove any existing bullet formatting
        for buNone in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone'):
            pPr.remove(buNone)

        # Add bullet character (solid circle) - RULE A
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        buChar.set('char', '\u2022')  # Solid circle bullet

        # Set bullet font size to 85% of text size (Rule A: 80-90%)
        # Text is 14pt (T3), so bullet should be ~12pt
        buSzPct = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct')
        buSzPct.set('val', '85000')  # 85% in percentage units (85000 = 85%)

        # Set bullet color to Primary Green - RULE A
        buClr = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buClr')
        srgbClr = etree.SubElement(buClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgbClr.set('val', bullet_config['bullet_style']['color'].lstrip('#').upper())

        # RULE B + RULE C: Hanging indent with text gutter
        # The gap between bullet and text is controlled by the difference between
        # left_indent and first_line_indent
        # Access via pPr XML element to set indent values (etree already imported above)
        pPr_elem = paragraph._element.get_or_add_pPr()

        if level == 0:
            # RULE C: 12px gap between bullet and text start (at 96 DPI final render)
            # RULE B: Text wraps to align with first line, not bullet
            # Set marL (left margin) and indent (first line indent) in EMUs
            # Note: px_to_inches() uses 144 DPI, so multiply by 1.5 to get correct 96 DPI render
            left_indent_emu = int(px_to_inches(48))  # 48px at 144 DPI = 32px at 96 DPI
            first_line_indent_emu = int(px_to_inches(-30))  # -30px at 144 DPI = -20px at 96 DPI

            pPr_elem.set('marL', str(left_indent_emu))
            pPr_elem.set('indent', str(first_line_indent_emu))
            # Effective text gutter = 48 - 30 = 18px at 144 DPI = 12px at 96 DPI

        elif level == 1:
            left_indent_emu = int(px_to_inches(93))  # 48 + 45px spacing
            first_line_indent_emu = int(px_to_inches(-30))
            pPr_elem.set('marL', str(left_indent_emu))
            pPr_elem.set('indent', str(first_line_indent_emu))

        elif level == 2:
            left_indent_emu = int(px_to_inches(138))  # 93 + 45px spacing
            first_line_indent_emu = int(px_to_inches(-30))
            pPr_elem.set('marL', str(left_indent_emu))
            pPr_elem.set('indent', str(first_line_indent_emu))

    def _get_dynamic_chart_typography(self, categories: list, series_count: int = 1) -> dict:
        """
        GLOBAL SKILL 4 ENHANCEMENT: Dynamic Typography for Chart Axes

        Automatically adjusts font sizes based on label density to ensure
        charts remain clean and readable.

        Rules:
        - Standard density (≤5 categories): Use normal sizing (9pt axes, 9pt legend)
        - Medium density (6-8 categories): Reduce to subordinate sizing (8pt)
        - High density (>8 categories): Reduce to minimum sizing (7pt)
        - Also considers average label length for additional adjustment

        Args:
            categories: List of category labels
            series_count: Number of data series (affects legend sizing)

        Returns:
            Dictionary with font sizes for axis_labels, legend, and data_labels
        """
        num_categories = len(categories)
        avg_label_length = sum(len(str(cat)) for cat in categories) / num_categories if num_categories > 0 else 0

        # Calculate density score (higher = more crowded)
        density_score = num_categories + (avg_label_length / 3)  # Weight label length

        # Standard sizing
        if density_score <= 6:
            return {
                'axis_labels': Pt(9),      # T4.5 standard
                'legend': Pt(9),           # T4.5 standard
                'data_labels': Pt(9),      # T4 standard
                'description': 'standard'
            }
        # Medium density - subordinate sizing
        elif density_score <= 10:
            return {
                'axis_labels': Pt(8),      # Reduced
                'legend': Pt(8),           # Reduced
                'data_labels': Pt(8),      # Reduced
                'description': 'subordinate'
            }
        # High density - minimum sizing
        else:
            return {
                'axis_labels': Pt(7),      # Minimum readable
                'legend': Pt(7),           # Minimum readable
                'data_labels': Pt(7),      # Minimum readable
                'description': 'minimum'
            }

    def _create_column_chart(self, slide, chart_data: dict, left, top, width, height):
        """
        Create a column chart with professional styling.

        This function reads styling rules from config and applies them meticulously:
        - Set bar colors (default: Light Grey, highlight: Primary Green)
        - Remove chart border
        - Set gridline style (1pt horizontal, Gridline Light)
        - Apply data labels (Outside End position)
        - GLOBAL SKILL 4: Dynamic typography based on label density

        Args:
            slide: PowerPoint slide object
            chart_data: Dictionary with 'categories' and 'series' keys
            left, top, width, height: Chart position and size

        Returns:
            Chart shape object
        """
        # Get chart specifications from config
        chart_spec = self.config['charts']['column']

        # Prepare chart data
        chart_data_obj = CategoryChartData()

        categories = chart_data.get('categories', ['Category 1', 'Category 2', 'Category 3'])
        series_list = chart_data.get('series', [{'name': 'Series 1', 'values': [10, 20, 15]}])

        chart_data_obj.categories = categories

        for series in series_list:
            chart_data_obj.add_series(series.get('name', 'Data'), series.get('values', []))

        # Add chart to slide
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            left, top, width, height,
            chart_data_obj
        )

        chart = chart_shape.chart

        # APPLY PROFESSIONAL STYLING

        # GLOBAL SKILL 4: Get dynamic typography based on label density
        dynamic_typography = self._get_dynamic_chart_typography(categories, len(series_list))

        # 2. Get plot for configuration
        plot = chart.plots[0]
        plot.has_data_labels = True

        # 3. GLOBAL SKILL 1: STORY-DRIVEN COLOR ENGINE
        # Apply intelligent color logic based on the story being told
        num_series = len(chart.series)
        story_colors = self.config.get('story_driven_colors', {})
        comparison_palette_rgb = story_colors.get('comparison_palette_rgb', [
            (2, 86, 69), (81, 123, 112), (81, 163, 163), (162, 218, 217)
        ])

        # Determine color strategy
        if isinstance(chart_data, dict) and chart_data.get('color_mode') == 'comparison':
            # COMPARISON MODE: Multiple series, each gets different color from palette
            color_mode = 'comparison'
        elif num_series == 1 and len(chart.series[0].points) > 1:
            # CATEGORY COMPARISON MODE: Single series with multiple categories
            # Each category/bar gets different color from palette
            color_mode = 'category_comparison'
        else:
            # HIGHLIGHT MODE (default): Grey + single highlight color
            color_mode = 'highlight'
            highlight_idx = chart_data.get('highlight_index', len(chart.series) - 1) if isinstance(chart_data, dict) else len(chart.series) - 1
            default_grey_rgb = story_colors.get('default_neutral_rgb', (211, 211, 211))
            highlight_rgb = story_colors.get('highlight_rgb', (20, 123, 88))  # Primary Green

        # LEGEND CONFIGURATION: Only show when it adds information
        # Hide legend for category_comparison mode - X-axis labels already identify each bar
        if color_mode == 'category_comparison':
            chart.has_legend = False  # Redundant - categories labeled on X-axis
        else:
            # Show legend for comparison (multi-series) or highlight mode
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            legend_spec = get_typography('T4.5')
            chart.legend.font.size = dynamic_typography['legend']
            chart.legend.font.name = 'Arial'
            chart.legend.font.color.rgb = RGBColor(*legend_spec['color_rgb'])

        # GLOBAL SKILL 2: LEGEND-TO-DATA INTEGRITY
        # Ensure legend colors match data colors exactly
        for idx, series in enumerate(chart.series):
            if color_mode == 'comparison':
                # Multi-series comparison: each series gets different color
                series_color_rgb = comparison_palette_rgb[idx % len(comparison_palette_rgb)]

                # Set color on SERIES FORMAT (for legend)
                try:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(*series_color_rgb)
                    series.format.line.fill.background()  # Borderless
                except:
                    pass

                # Apply same color to all points in this series
                for point in series.points:
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor(*series_color_rgb)
                    point.format.line.fill.background()

            elif color_mode == 'category_comparison':
                # Single-series category comparison: each point/category gets different color
                # Don't set series-level color (would show wrong legend)
                for point_idx, point in enumerate(series.points):
                    point_color_rgb = comparison_palette_rgb[point_idx % len(comparison_palette_rgb)]
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor(*point_color_rgb)
                    point.format.line.fill.background()

            else:
                # Highlight mode: grey for context, color for story
                is_highlight = (idx == highlight_idx)
                series_color_rgb = highlight_rgb if is_highlight else default_grey_rgb

                # Set color on SERIES FORMAT (for legend)
                try:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(*series_color_rgb)
                    series.format.line.fill.background()  # Borderless
                except:
                    pass

                # Apply same color to all points in this series
                for point in series.points:
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor(*series_color_rgb)
                    point.format.line.fill.background()

        # 4. Set gap width (50% of bar width)
        plot.gap_width = chart_spec['bars']['gap_width_percent']

        # 5. Configure value axis (Y-axis)
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(*chart_spec['gridlines']['horizontal']['color_rgb'])
        value_axis.major_gridlines.format.line.width = chart_spec['gridlines']['horizontal']['line_thickness']

        # Remove Y-axis line
        value_axis.format.line.fill.background()

        # GLOBAL SKILL 4: Dynamic axis label sizing based on density
        # Chart elements must be visually subordinate to main content
        axis_spec = get_typography('T4.5')
        axis_color_rgb = axis_spec['color_rgb']  # Axis Grey

        value_axis.tick_labels.font.size = dynamic_typography['axis_labels']  # Dynamic size
        value_axis.tick_labels.font.name = 'Arial'
        value_axis.tick_labels.font.color.rgb = RGBColor(*axis_color_rgb)

        # 6. Configure category axis (X-axis)
        category_axis = chart.category_axis
        category_axis.format.line.color.rgb = RGBColor(*chart_spec['axes']['x_axis']['color_rgb'])
        category_axis.format.line.width = chart_spec['axes']['x_axis']['line_thickness']

        # GLOBAL SKILL 4: Dynamic category axis label sizing
        category_axis.tick_labels.font.size = dynamic_typography['axis_labels']  # Dynamic size
        category_axis.tick_labels.font.name = 'Arial'
        category_axis.tick_labels.font.color.rgb = RGBColor(*axis_color_rgb)

        # 7. Apply data labels (T4 typography, Outside End position)
        # GLOBAL SKILL 4: Dynamic data label sizing based on density
        data_labels = plot.data_labels
        data_labels.font.size = dynamic_typography['data_labels']  # Dynamic size
        data_labels.font.name = 'Arial'
        data_labels.font.color.rgb = RGBColor(*self.colors['primary_body_text_rgb'])
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        # DATA FORMATTING CONTROL: No decimal places for cleaner presentation
        # Integer display (e.g., "250" not "250.0") for clarity
        data_labels.number_format = '#,##0'  # No decimals, thousands separator

        # IMPROVEMENT 4: Add vertical padding (5-8px) so labels "float" above bars
        # Note: python-pptx doesn't expose offset directly, but OUTSIDE_END positioning
        # automatically adds some padding. The 10pt smaller font also helps with visual spacing.

        return chart_shape

    def _add_footer(self, slide, slide_content):
        """
        Add footer with source attribution.

        PRINCIPLE 1 - Zoned Integrity: Footer is confined to Footer zone,
        cannot be violated by Content zone elements.

        PRINCIPLE 5 - Consistent Visual Anchoring: Footer text vertically
        centered within Footer zone.

        Args:
            slide: PowerPoint slide object
            slide_content: Slide content with optional source
        """
        footer_region = get_region('footer')
        footer_spec = get_typography('T5')

        if hasattr(slide_content, 'source') and slide_content.source:
            # PRINCIPLE 1: Footer constrained to Footer zone bounds
            footer_box = slide.shapes.add_textbox(
                footer_region['bounds']['x1'],
                footer_region['bounds']['y1'],
                px_to_inches(footer_region['bounds']['x2_px'] - footer_region['bounds']['x1_px']),
                footer_region['height']
            )

            footer_frame = footer_box.text_frame
            footer_frame.text = f"Source: {slide_content.source}"
            footer_frame.word_wrap = False
            footer_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # PRINCIPLE 5: Vertical center
            footer_frame.margin_left = Inches(0)
            footer_frame.margin_right = Inches(0)

            p = footer_frame.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            p.font.name = 'Arial'
            p.font.size = footer_spec['font_size']  # T5: 9pt
            p.font.color.rgb = RGBColor(*footer_spec['color_rgb'])

    def _add_closing_slide(self, prs: Presentation, request: PresentationRequest):
        """Add closing slide with professional styling."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Set background to Primary Green (section divider style)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.colors['primary_green_rgb'])

        # Get section divider typography
        divider_spec = self.config['typography_special']['section_divider_title']

        # Add closing text (vertically centered, left-aligned)
        closing_box = slide.shapes.add_textbox(
            divider_spec['position']['x'],
            divider_spec['position']['y'],
            px_to_inches(1728),
            px_to_inches(108)
        )

        closing_frame = closing_box.text_frame
        closing_frame.text = "Thank you"
        closing_frame.word_wrap = True
        closing_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Format closing text
        p = closing_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = 'Arial'
        p.font.size = divider_spec['font_size']
        p.font.bold = True
        p.font.color.rgb = RGBColor(*divider_spec['color_rgb'])

        # Add subtitle if author provided
        if request.author or request.company:
            subtitle_text = " | ".join(filter(None, [request.author, request.company]))

            subtitle_box = slide.shapes.add_textbox(
                px_to_inches(96),
                divider_spec['position']['y'] + px_to_inches(120),
                px_to_inches(1728),
                px_to_inches(60)
            )

            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            p = subtitle_frame.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.bold = False
            p.font.color.rgb = RGBColor(*self.colors['white_rgb'])

    def _add_chart_title(self, slide, title_text: str, left, top, width, height):
        """
        CHART TITLE SKILL: Add mandatory descriptive title using T2 Typography (20pt Bold).

        This is a REQUIRED element for all charts in consulting presentations.
        The chart title is considered part of the chart element for vertical centering.

        PRINCIPLE 4 - Hierarchical Titling: Chart title is secondary descriptive
        title (T2) explaining "what it is" (e.g., "Market Size Growth (EUR B)"),
        positioned below the main slide title (T1) which explains "so what".

        PRINCIPLE 5 - Consistent Visual Anchoring: Chart title vertically
        centered within its allocated space for visual balance.

        Args:
            slide: PowerPoint slide object
            title_text: Chart title text (T2: descriptive, e.g., "Revenue by Channel (EUR M)")
            left, top, width, height: Position and size in inches
        """
        title_spec = get_typography('T2')  # Now 20pt Bold (upgraded from 18pt)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # PRINCIPLE 5: Vertical center
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)

        p = title_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = 'Arial'
        p.font.size = title_spec['font_size']  # T2: 18pt
        p.font.bold = True
        p.font.color.rgb = RGBColor(*title_spec['color_rgb'])  # Body Text color

    def _add_chart_source(self, slide, source_text: str, chart_area: dict):
        """
        FINAL PROFESSIONAL STANDARD 4: Chart Source Annotation Skill.

        Place data source in bottom-left corner of chart area (not slide footer).
        This directly attributes the data to its source - a standard practice in
        consulting presentations.

        Args:
            slide: PowerPoint slide object
            source_text: Source text (e.g., "Source: Company data", "Source: Analysis")
            chart_area: Chart area dictionary with position and size
        """
        source_spec = get_typography('T5')  # T5: 9pt, Axis Grey

        # Helper to convert to inches float
        def to_inches_float(val):
            if hasattr(val, 'inches'):
                return val.inches
            else:
                return val / 914400  # Assume EMU

        # Position in bottom-left corner of chart area
        # Place 10px from left edge and 10px from bottom edge
        source_height_inches = 0.25  # Small textbox height
        source_width_inches = 3.0  # Reasonable width for source text
        margin_inches = 10 / 144  # 10px at 144 DPI

        # Calculate position using raw inch values, then convert to Inches objects
        left_inches = to_inches_float(chart_area['x1']) + margin_inches
        top_inches = to_inches_float(chart_area['y1']) + to_inches_float(chart_area['height']) - source_height_inches - margin_inches

        source_box = slide.shapes.add_textbox(
            Inches(left_inches),
            Inches(top_inches),
            Inches(source_width_inches),
            Inches(source_height_inches)
        )
        source_frame = source_box.text_frame
        source_frame.text = source_text
        source_frame.word_wrap = False

        p = source_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = 'Arial'
        p.font.size = source_spec['font_size']  # T5: 9pt
        p.font.bold = False
        p.font.color.rgb = RGBColor(*source_spec['color_rgb'])  # Axis Grey

    def _create_waterfall_chart(self, slide, chart_data: dict, left, top, width, height):
        """
        ENHANCEMENT 3: Create a waterfall chart with professional styling.

        Waterfall charts show incremental changes from a starting value to ending value,
        useful for profit drivers, cost breakdowns, etc.

        Args:
            slide: PowerPoint slide object
            chart_data: Dictionary with 'categories', 'values', and optional 'types'
            left, top, width, height: Chart position and size

        Returns:
            Chart shape object
        """
        # Get chart specifications
        chart_spec = self.config['charts'].get('waterfall', {})

        # Prepare chart data
        chart_data_obj = CategoryChartData()

        categories = chart_data.get('categories', [])
        values = chart_data.get('values', [])
        types = chart_data.get('types', [])  # 'start', 'increase', 'decrease', 'end'

        chart_data_obj.categories = categories
        chart_data_obj.add_series('Waterfall', values)

        # Add chart to slide (using column chart as base)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            left, top, width, height,
            chart_data_obj
        )

        chart = chart_shape.chart

        # Configure chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.gap_width = 50

        # GLOBAL SKILL 4: Get dynamic typography based on label density
        dynamic_typography = self._get_dynamic_chart_typography(categories, 1)

        # GLOBAL SKILL 1: Apply Story-Driven Color Engine to waterfall
        # Waterfall uses distinct colors: Primary Green for baseline, Accent Blue for increases, Red for decreases
        primary_green_rgb = (20, 123, 88)   # Start/End - baseline values
        accent_blue_rgb = (0, 94, 184)       # Increase - positive drivers
        negative_red_rgb = (230, 81, 102)    # Decrease - negative drivers

        series = chart.series[0]
        for idx, point in enumerate(series.points):
            point_type = types[idx] if idx < len(types) else 'increase'

            if point_type in ['start', 'end']:
                # Start/End: Primary Green - these are baseline totals
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor(*primary_green_rgb)
            elif point_type == 'increase':
                # Increase: Accent Blue - clearly distinguishes positive drivers from baseline
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor(*accent_blue_rgb)
            else:  # decrease
                # Decrease: Negative Red - shows negative impact
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor(*negative_red_rgb)

            # FINAL PROFESSIONAL STANDARD 3: Borderless
            point.format.line.fill.background()

        # Configure axes with GLOBAL SKILL 4: Dynamic typography
        axis_spec = get_typography('T4.5')
        axis_color_rgb = axis_spec['color_rgb']

        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(224, 224, 224)
        value_axis.major_gridlines.format.line.width = Pt(1)
        value_axis.format.line.fill.background()
        value_axis.tick_labels.font.size = dynamic_typography['axis_labels']  # Dynamic size
        value_axis.tick_labels.font.name = 'Arial'
        value_axis.tick_labels.font.color.rgb = RGBColor(*axis_color_rgb)

        category_axis = chart.category_axis
        category_axis.format.line.color.rgb = RGBColor(169, 169, 169)
        category_axis.format.line.width = Pt(1.5)
        category_axis.tick_labels.font.size = dynamic_typography['axis_labels']  # Dynamic size
        category_axis.tick_labels.font.name = 'Arial'
        category_axis.tick_labels.font.color.rgb = RGBColor(*axis_color_rgb)

        # Data labels - GLOBAL SKILL 4: Dynamic typography
        data_labels = plot.data_labels
        data_labels.font.size = dynamic_typography['data_labels']  # Dynamic size
        data_labels.font.name = 'Arial'
        data_labels.font.color.rgb = RGBColor(*self.colors['primary_body_text_rgb'])
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.number_format = '#,##0'  # No decimals, thousands separator

        return chart_shape

    def _create_matrix_chart(self, slide, chart_data: dict, left, top, width, height):
        """
        ENHANCEMENT 3: Create Matrix chart (2x2 bubble chart with reversed X-axis).

        The Matrix is a strategic planning tool showing:
        - X-axis: Relative Market Share (HIGH to LOW - reversed)
        - Y-axis: Market Growth (LOW to HIGH)
        - Quadrants: Stars, Question Marks, Cash Cows, Dogs
        - Bubble size: Revenue or strategic importance

        Args:
            slide: PowerPoint slide object
            chart_data: Dictionary with 'bubbles' list containing x, y, size, label
            left, top, width, height: Chart position and size

        Returns:
            Chart shape object
        """
        # Get chart specifications
        chart_spec = self.config['charts'].get('matrix', {})

        # Prepare bubble chart data
        chart_data_obj = BubbleChartData()

        bubbles = chart_data.get('bubbles', [])

        # Add each bubble as a separate series (for individual coloring)
        for bubble in bubbles:
            series = chart_data_obj.add_series(bubble.get('label', 'Item'))
            series.add_data_point(
                bubble.get('x', 1.0),
                bubble.get('y', 1.0),
                bubble.get('size', 50)
            )

        # Add chart to slide
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BUBBLE,
            left, top, width, height,
            chart_data_obj
        )

        chart = chart_shape.chart

        # Configure chart
        chart.has_legend = False

        # Reverse X-axis (high market share on left)
        category_axis = chart.category_axis
        category_axis.reverse_order = True
        category_axis.has_title = True
        category_axis.axis_title.text_frame.text = "Relative Market Share"
        category_axis.format.line.color.rgb = RGBColor(169, 169, 169)
        category_axis.format.line.width = Pt(1.5)

        # Configure Y-axis
        value_axis = chart.value_axis
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = "Market Growth"
        value_axis.format.line.color.rgb = RGBColor(169, 169, 169)
        value_axis.format.line.width = Pt(1.5)

        # Style bubbles with green-based sequential colors
        colors = [(2, 86, 69), (81, 123, 112), (81, 163, 163), (162, 218, 217)]  # Green sequential palette
        for idx, series in enumerate(chart.series):
            color = colors[idx % len(colors)]
            for point in series.points:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor(*color)
                point.format.line.color.rgb = RGBColor(74, 74, 74)
                point.format.line.width = Pt(1)

        return chart_shape

    def _add_chart_annotations(self, slide, chart_shape, annotations: list, chart_area: dict):
        """
        ENHANCEMENT 2: Add advanced chart annotations.

        Supports:
        - CAGR arrows with growth rate labels
        - Difference lines between data points
        - Callout boxes pointing to specific data points

        Args:
            slide: PowerPoint slide object
            chart_shape: Chart shape to annotate
            annotations: List of annotation dictionaries
            chart_area: Chart area bounds for positioning
        """
        for annotation in annotations:
            ann_type = annotation.get('type')

            if ann_type == 'cagr_arrow':
                self._add_cagr_arrow(slide, chart_shape, annotation, chart_area)
            elif ann_type == 'leader_line':
                # GLOBAL SKILL 3: Leader line for specific data points
                self._add_leader_line(slide, annotation, chart_area)
            elif ann_type == 'difference_line':
                self._add_difference_line(slide, chart_shape, annotation, chart_area)
            elif ann_type == 'callout':
                self._add_callout(slide, annotation, chart_area)

    def _add_cagr_arrow(self, slide, chart_shape, annotation: dict, chart_area: dict):
        """
        GLOBAL SKILL 3: Add CAGR arrow showing growth from start to end point.

        FULLY DYNAMIC AND DATA-DRIVEN IMPLEMENTATION:
        1. Get Anchor Coordinates: Programmatically determine exact (X, Y) of bar tops
        2. Draw Curved Line: Quadratic Bézier curve with vertical lift offset
        3. Calculate Apex: Determine highest point of the curve
        4. Position Label: Place text precisely above the apex, horizontally centered

        All positions are derived procedurally from chart data. No hardcoded coordinates.

        Args:
            slide: PowerPoint slide object
            chart_shape: Chart shape
            annotation: Dict with 'series_index', 'from_category', 'to_category', 'label'
            chart_area: Chart area bounds (Inches objects)
        """
        from pptx.util import Emu
        from pptx.oxml.ns import nsmap
        from pptx.oxml import parse_xml
        from lxml import etree

        chart = chart_shape.chart

        # Helper to convert to inches float
        def to_inches_float(val):
            if hasattr(val, 'inches'):
                return val.inches
            elif hasattr(val, 'pt'):
                return val.inches
            else:
                return val / 914400

        # =========================================================================
        # STEP 1: GET ANCHOR COORDINATES FROM CHART DATA
        # =========================================================================
        chart_x1 = to_inches_float(chart_area['x1'])
        chart_y1 = to_inches_float(chart_area['y1'])
        chart_width = to_inches_float(chart_area['width'])
        chart_height = to_inches_float(chart_area['height'])

        # Get series and category indices
        series_idx = annotation.get('series_index', 0)
        from_cat_idx = annotation.get('from_category', 0)
        to_cat_idx = annotation.get('to_category', -1)

        # Get chart data
        series = chart.series[series_idx]
        num_series = len(chart.series)
        num_categories = len(series.points)
        if to_cat_idx < 0:
            to_cat_idx = num_categories + to_cat_idx

        # Get data values
        series_values = list(series.values)
        from_value = series_values[from_cat_idx]
        to_value = series_values[to_cat_idx]

        # Get Y-axis range (auto-calculated or explicit)
        value_axis = chart.value_axis
        all_values = []
        for s in chart.series:
            all_values.extend(list(s.values))
        data_min = min(all_values)
        data_max = max(all_values)

        y_min = value_axis.minimum_scale if value_axis.minimum_scale is not None else 0
        y_max = value_axis.maximum_scale if value_axis.maximum_scale is not None else data_max * 1.1
        y_range = y_max - y_min
        if y_range <= 0:
            y_range = 1

        # Calculate plot area within chart (accounting for axes, labels, legend)
        # Estimate plot area margins (PowerPoint typically uses ~15% for axes/labels)
        plot_left_margin = 0.12  # 12% for Y-axis labels
        plot_right_margin = 0.02  # 2% right padding
        plot_top_margin = 0.08  # 8% for title/padding
        plot_bottom_margin = 0.15  # 15% for X-axis labels and legend

        plot_x1 = chart_x1 + (plot_left_margin * chart_width)
        plot_width = chart_width * (1 - plot_left_margin - plot_right_margin)
        plot_y1 = chart_y1 + (plot_top_margin * chart_height)
        plot_height = chart_height * (1 - plot_top_margin - plot_bottom_margin)

        # Calculate bar positions within plot area
        # For clustered bar charts: bars are grouped by category
        gap_width_ratio = 0.5  # 50% gap (matches chart config)
        category_width = plot_width / num_categories
        bar_group_width = category_width / (1 + gap_width_ratio)
        single_bar_width = bar_group_width / num_series
        gap_between_categories = category_width - bar_group_width

        # Calculate X position: center of the specific bar within its category
        def get_bar_center_x(cat_idx, ser_idx):
            category_start = plot_x1 + (cat_idx * category_width) + (gap_between_categories / 2)
            bar_start = category_start + (ser_idx * single_bar_width)
            bar_center = bar_start + (single_bar_width / 2)
            return bar_center

        # Calculate Y position: top of bar based on data value
        def get_bar_top_y(value):
            value_ratio = (value - y_min) / y_range
            # Y increases downward in PowerPoint, so invert
            bar_top_y = plot_y1 + plot_height - (value_ratio * plot_height)
            return bar_top_y

        # Get anchor points (top-center of start and end bars)
        from_x = get_bar_center_x(from_cat_idx, series_idx)
        from_y = get_bar_top_y(from_value)
        to_x = get_bar_center_x(to_cat_idx, series_idx)
        to_y = get_bar_top_y(to_value)

        # =========================================================================
        # STEP 2: DRAW QUADRATIC BÉZIER CURVE (DATA-DRIVEN ARC HEIGHT)
        # =========================================================================
        # NEW LOGIC: Arc height is determined by the highest obstacle between start and end

        # STEP 2a: Identify the highest obstacle between start and end categories
        # Find the maximum value of any bar (across all series) between from_cat and to_cat
        min_cat = min(from_cat_idx, to_cat_idx)
        max_cat = max(from_cat_idx, to_cat_idx)

        highest_obstacle_value = 0
        for cat_idx in range(min_cat, max_cat + 1):
            for s in chart.series:
                s_values = list(s.values)
                if cat_idx < len(s_values):
                    highest_obstacle_value = max(highest_obstacle_value, s_values[cat_idx])

        # STEP 2b: Calculate the arc's peak with clearance factor
        # The peak should be just above the highest obstacle
        clearance_px = 30  # 30 pixels clearance above highest bar
        clearance_inches = clearance_px / 144  # Convert from 144 DPI reference

        # Get the Y position of the highest obstacle
        highest_obstacle_y = get_bar_top_y(highest_obstacle_value)

        # The control point (and thus the arc peak) should be clearance_inches above the highest obstacle
        # For a quadratic Bézier, the curve reaches approximately 75% of the way to the control point
        # So we need to set control_y such that the actual peak clears by clearance_inches
        # Adjustment factor: control point needs to be ~1.33x the desired peak offset
        control_y = highest_obstacle_y - (clearance_inches * 1.33)

        mid_x = (from_x + to_x) / 2

        # Generate points along the Bézier curve for the freeform shape
        # Quadratic Bézier: B(t) = (1-t)²P0 + 2(1-t)tP1 + t²P2
        def bezier_point(t, p0, p1, p2):
            return ((1 - t) ** 2) * p0 + 2 * (1 - t) * t * p1 + (t ** 2) * p2

        num_segments = 20
        curve_points = []
        for i in range(num_segments + 1):
            t = i / num_segments
            bx = bezier_point(t, from_x, mid_x, to_x)
            by = bezier_point(t, from_y, control_y, to_y)
            curve_points.append((bx, by))

        # Create freeform shape for the curved line
        # Use a series of line segments to approximate the curve
        builder = slide.shapes.build_freeform(
            Inches(curve_points[0][0]),
            Inches(curve_points[0][1])
        )
        for px, py in curve_points[1:]:
            builder.add_line_segments([(Inches(px), Inches(py))], close=False)

        curve_shape = builder.convert_to_shape()
        curve_shape.line.color.rgb = RGBColor(169, 169, 169)  # Axis Grey
        curve_shape.line.width = Pt(1)  # Thin line
        curve_shape.fill.background()  # No fill

        # =========================================================================
        # STEP 3: CALCULATE THE APEX OF THE CURVE
        # =========================================================================
        # For a quadratic Bézier, the apex (highest point) occurs at t where dy/dt = 0
        # For our purposes, we can find it by checking the curve points
        apex_y = min(p[1] for p in curve_points)
        apex_idx = [p[1] for p in curve_points].index(apex_y)
        apex_x = curve_points[apex_idx][0]

        # =========================================================================
        # STEP 4: POSITION THE LABEL ABOVE THE APEX
        # =========================================================================
        label_text = annotation.get('label', 'CAGR')

        # Position label above apex, horizontally centered
        label_width = 0.9  # Inches
        label_height = 0.25  # Inches
        label_offset_above = 0.08  # Inches above apex

        label_x = apex_x - (label_width / 2)  # Center horizontally on apex
        label_y = apex_y - label_height - label_offset_above  # Above the apex

        label_box = slide.shapes.add_textbox(
            Inches(label_x),
            Inches(label_y),
            Inches(label_width),
            Inches(label_height)
        )
        label_frame = label_box.text_frame
        label_frame.word_wrap = False
        label_frame.text = label_text
        p = label_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.size = Pt(10)
        p.font.bold = False
        p.font.color.rgb = RGBColor(74, 74, 74)  # Body text grey

    def _add_leader_line(self, slide, annotation: dict, chart_area: dict):
        """
        GLOBAL SKILL 3: Add leader line from data point to label.

        Visual Annotation System: Thin, straight line extending from data point
        to clean white space where text label is placed. Never floating freely.
        For specific data points (e.g., '10 Million Users').

        Args:
            slide: PowerPoint slide object
            annotation: Dict with 'x', 'y', 'text', 'direction' ('up', 'down', 'left', 'right')
            chart_area: Chart area bounds (Inches objects)
        """
        # Helper to convert to inches float
        def to_inches_float(val):
            if hasattr(val, 'inches'):
                return val.inches
            else:
                return val / 914400

        chart_x1 = to_inches_float(chart_area['x1'])
        chart_y1 = to_inches_float(chart_area['y1'])
        chart_width = to_inches_float(chart_area['width'])
        chart_height = to_inches_float(chart_area['height'])

        # Data point position
        point_x_inches = chart_x1 + (annotation.get('x', 0.5) * chart_width)
        point_y_inches = chart_y1 + ((1 - annotation.get('y', 0.5)) * chart_height)

        # Leader line direction
        direction = annotation.get('direction', 'up')
        line_length = annotation.get('line_length', 40)  # pixels at 144 DPI

        # Calculate end point based on direction
        if direction == 'up':
            end_x_inches = point_x_inches
            end_y_inches = point_y_inches - (line_length / 144)
            label_x_inches = end_x_inches - (50 / 144)
            label_y_inches = end_y_inches - (25 / 144)
        elif direction == 'down':
            end_x_inches = point_x_inches
            end_y_inches = point_y_inches + (line_length / 144)
            label_x_inches = end_x_inches - (50 / 144)
            label_y_inches = end_y_inches + (5 / 144)
        elif direction == 'right':
            end_x_inches = point_x_inches + (line_length / 144)
            end_y_inches = point_y_inches
            label_x_inches = end_x_inches + (5 / 144)
            label_y_inches = end_y_inches - (10 / 144)
        else:  # left
            end_x_inches = point_x_inches - (line_length / 144)
            end_y_inches = point_y_inches
            label_x_inches = end_x_inches - (105 / 144)
            label_y_inches = end_y_inches - (10 / 144)

        # Add leader line - thin and grey
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(point_x_inches), Inches(point_y_inches),
            Inches(end_x_inches), Inches(end_y_inches)
        )
        line.line.color.rgb = RGBColor(169, 169, 169)  # Axis Grey
        line.line.width = Pt(0.75)  # Very thin

        # Add text label at end of leader line
        label_text = annotation.get('text', '')
        label_box = slide.shapes.add_textbox(
            Inches(label_x_inches),
            Inches(label_y_inches),
            Inches(100 / 144),
            Inches(20 / 144)
        )
        label_frame = label_box.text_frame
        label_frame.text = label_text
        label_frame.word_wrap = False
        p = label_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT if direction in ['right', 'up', 'down'] else PP_PARAGRAPH_ALIGNMENT.RIGHT
        p.font.size = Pt(10)
        p.font.bold = False
        p.font.color.rgb = RGBColor(74, 74, 74)

    def _add_difference_line(self, slide, chart_shape, annotation: dict, chart_area: dict):
        """
        GLOBAL SKILL 3: Add vertical difference line showing delta between two bars.

        VISUAL DESIGN AI IMPLEMENTATION:
        1. Extract bar geometry from chart data (positions in pixels at 144 DPI)
        2. Build JSON input for Visual Design AI
        3. Get precise placement coordinates (using AI or local fallback)
        4. Render the annotation using those coordinates

        The Visual Design AI follows these rules:
        - Line is placed at the gutter centerline (between bars)
        - Label is positioned with 10px padding, collision-aware
        - If label would overlap bar2, it moves to the left side

        Args:
            slide: PowerPoint slide object
            chart_shape: Chart shape
            annotation: Dict with 'series_index', 'from_category', 'to_category', 'label'
            chart_area: Chart area bounds (Inches objects)

        Returns:
            Tuple of (from_cat_idx, to_cat_idx) for data label suppression
        """
        chart = chart_shape.chart

        # Helper to convert to inches float
        def to_inches_float(val):
            if hasattr(val, 'inches'):
                return val.inches
            else:
                return val / 914400

        # =========================================================================
        # STEP 1: EXTRACT BAR GEOMETRY FROM CHART DATA
        # =========================================================================
        chart_x1 = to_inches_float(chart_area['x1'])
        chart_y1 = to_inches_float(chart_area['y1'])
        chart_width = to_inches_float(chart_area['width'])
        chart_height = to_inches_float(chart_area['height'])

        # Get series and category indices
        series_idx = annotation.get('series_index', 0)
        from_cat_idx = annotation.get('from_category', 0)
        to_cat_idx = annotation.get('to_category', 1)

        # Get chart data
        series = chart.series[series_idx]
        num_series = len(chart.series)
        num_categories = len(series.points)

        # Get data values
        series_values = list(series.values)
        from_value = series_values[from_cat_idx]
        to_value = series_values[to_cat_idx]

        # Get Y-axis range
        value_axis = chart.value_axis
        all_values = []
        for s in chart.series:
            all_values.extend(list(s.values))
        data_max = max(all_values)

        y_min = value_axis.minimum_scale if value_axis.minimum_scale is not None else 0
        y_max = value_axis.maximum_scale if value_axis.maximum_scale is not None else data_max * 1.1
        y_range = y_max - y_min
        if y_range <= 0:
            y_range = 1

        # Calculate plot area margins
        plot_left_margin = 0.12
        plot_right_margin = 0.02
        plot_top_margin = 0.08
        plot_bottom_margin = 0.15

        plot_x1 = chart_x1 + (plot_left_margin * chart_width)
        plot_width = chart_width * (1 - plot_left_margin - plot_right_margin)
        plot_y1 = chart_y1 + (plot_top_margin * chart_height)
        plot_height = chart_height * (1 - plot_top_margin - plot_bottom_margin)

        # Calculate bar positions
        gap_width_ratio = 0.5
        category_width = plot_width / num_categories
        bar_group_width = category_width / (1 + gap_width_ratio)
        single_bar_width = bar_group_width / num_series
        gap_between_categories = category_width - bar_group_width

        def get_bar_edges_inches(cat_idx, ser_idx):
            """Return (left_edge, right_edge) of a bar in inches."""
            category_start = plot_x1 + (cat_idx * category_width) + (gap_between_categories / 2)
            bar_start = category_start + (ser_idx * single_bar_width)
            bar_end = bar_start + single_bar_width
            return (bar_start, bar_end)

        def get_bar_top_y_inches(value):
            value_ratio = (value - y_min) / y_range
            bar_top_y = plot_y1 + plot_height - (value_ratio * plot_height)
            return bar_top_y

        # Get bar edges in inches
        from_bar_left, from_bar_right = get_bar_edges_inches(from_cat_idx, series_idx)
        to_bar_left, to_bar_right = get_bar_edges_inches(to_cat_idx, series_idx)
        from_bar_y = get_bar_top_y_inches(from_value)
        to_bar_y = get_bar_top_y_inches(to_value)

        # =========================================================================
        # STEP 2: BUILD JSON INPUT FOR VISUAL DESIGN AI (convert to pixels at 144 DPI)
        # =========================================================================
        DPI = 144  # Reference DPI for pixel coordinates

        bar1_geometry = {
            "left_edge_x": from_bar_left * DPI,
            "right_edge_x": from_bar_right * DPI,
            "top_edge_y": from_bar_y * DPI
        }
        bar2_geometry = {
            "left_edge_x": to_bar_left * DPI,
            "right_edge_x": to_bar_right * DPI,
            "top_edge_y": to_bar_y * DPI
        }
        label_text = annotation.get('label', 'Delta')

        # =========================================================================
        # STEP 3: GET PLACEMENT FROM VISUAL DESIGN AI (or local fallback)
        # =========================================================================
        # Use local Visual Design AI logic (same rules as LLM would apply)
        placement = self._visual_design_ai_place_difference_line(
            bar1_geometry, bar2_geometry, label_text
        )

        # =========================================================================
        # STEP 4: RENDER ANNOTATION USING AI-CALCULATED COORDINATES
        # =========================================================================
        # Convert pixel coordinates back to inches
        line_x = placement["line"]["position_x"] / DPI
        line_y1 = placement["line"]["start_y"] / DPI
        line_y2 = placement["line"]["end_y"] / DPI

        # Draw vertical dashed red line
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(line_x), Inches(line_y1),
            Inches(line_x), Inches(line_y2)
        )
        line.line.color.rgb = RGBColor(230, 81, 102)  # Negative Red
        line.line.width = Pt(1.5)
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # Deconstruct label into primary and secondary lines
        primary_line = label_text
        secondary_line = None

        if '\n' in label_text:
            parts = label_text.split('\n', 1)
            primary_line = parts[0].strip()
            secondary_line = parts[1].strip() if len(parts) > 1 else None
        elif '(' in label_text:
            paren_idx = label_text.find('(')
            primary_line = label_text[:paren_idx].strip()
            secondary_line = label_text[paren_idx:].strip()

        # Get label position from placement
        label_x = placement["label"]["position_x"] / DPI
        label_center_y = placement["label"]["vertical_center_y"] / DPI

        # Label dimensions
        label_width = 1.0  # Inches
        label_height = 0.5 if secondary_line else 0.3  # Inches
        label_y = label_center_y - (label_height / 2)

        # Create the text box
        label_box = slide.shapes.add_textbox(
            Inches(label_x),
            Inches(label_y),
            Inches(label_width),
            Inches(label_height)
        )
        label_frame = label_box.text_frame
        label_frame.word_wrap = False
        label_frame.auto_size = MSO_AUTO_SIZE.NONE

        # Style primary line (bold, 11pt)
        p1 = label_frame.paragraphs[0]
        p1.text = primary_line
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(230, 81, 102)  # Negative Red
        p1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Add secondary line if present (regular, 10pt)
        if secondary_line:
            p2 = label_frame.add_paragraph()
            p2.text = secondary_line
            p2.font.size = Pt(10)
            p2.font.bold = False
            p2.font.color.rgb = RGBColor(230, 81, 102)  # Negative Red
            p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Return the annotated category indices for data label suppression
        return (from_cat_idx, to_cat_idx)

    def _visual_design_ai_place_difference_line(
        self,
        bar1: dict,
        bar2: dict,
        label_text: str
    ) -> dict:
        """
        Visual Design AI: Calculate precise placement for difference line annotation.

        This method implements the Visual Design AI rules for annotation placement:
        1. Calculate the Gutter: line.position_x = midpoint of empty space between bars
        2. Set Line Height: match bar top edges
        3. Position Label: 10px to the right, vertically centered
        4. Collision Avoidance: if label overlaps bar2, move to left side

        Args:
            bar1: Dict with 'left_edge_x', 'right_edge_x', 'top_edge_y' (pixels)
            bar2: Dict with 'left_edge_x', 'right_edge_x', 'top_edge_y' (pixels)
            label_text: The annotation label text

        Returns:
            Dict with 'line' and 'label' placement data
        """
        # RULE 1: Calculate the Gutter Centerline
        gutter_centerline = (bar1["right_edge_x"] + bar2["left_edge_x"]) / 2

        # RULE 2: Set Line Height (bar tops)
        line_start_y = bar1["top_edge_y"]
        line_end_y = bar2["top_edge_y"]

        # RULE 3: Position the Label
        vertical_center_y = (line_start_y + line_end_y) / 2
        label_padding = 10  # 10 pixels
        label_width_estimate = 144  # ~1 inch at 144 DPI

        # Default: place to the right
        label_x = gutter_centerline + label_padding
        placement_side = "right"

        # RULE 4: Collision Avoidance
        if label_x + label_width_estimate > bar2["left_edge_x"]:
            # Collision detected - move to left
            label_x = gutter_centerline - label_padding - label_width_estimate
            placement_side = "left"

        return {
            "line": {
                "type": "dashed",
                "color": "Negative Red",
                "start_y": line_start_y,
                "end_y": line_end_y,
                "position_x": gutter_centerline
            },
            "label": {
                "text": label_text,
                "position_x": label_x,
                "vertical_center_y": vertical_center_y,
                "placement_side": placement_side
            }
        }

    def _add_callout(self, slide, annotation: dict, chart_area: dict):
        """
        Add callout box pointing to a specific data point.

        Args:
            slide: PowerPoint slide object
            annotation: Dict with 'x', 'y', 'text', 'position'
            chart_area: Chart area bounds (Inches objects)
        """
        # Helper to convert to inches float
        def to_inches_float(val):
            if hasattr(val, 'inches'):
                return val.inches
            else:
                return val / 914400

        chart_x1 = to_inches_float(chart_area['x1'])
        chart_y1 = to_inches_float(chart_area['y1'])
        chart_width = to_inches_float(chart_area['width'])
        chart_height = to_inches_float(chart_area['height'])

        # Calculate callout position
        point_x_inches = chart_x1 + (annotation.get('x', 0.5) * chart_width)
        point_y_inches = chart_y1 + ((1 - annotation.get('y', 0.5)) * chart_height)

        # Position callout based on preference
        position = annotation.get('position', 'above')
        if position == 'above':
            callout_x_inches = point_x_inches - (60 / 144)
            callout_y_inches = point_y_inches - (50 / 144)
        elif position == 'below':
            callout_x_inches = point_x_inches - (60 / 144)
            callout_y_inches = point_y_inches + (10 / 144)
        elif position == 'right':
            callout_x_inches = point_x_inches + (10 / 144)
            callout_y_inches = point_y_inches - (20 / 144)
        else:  # left
            callout_x_inches = point_x_inches - (130 / 144)
            callout_y_inches = point_y_inches - (20 / 144)

        # Add callout box
        callout_box = slide.shapes.add_textbox(
            Inches(callout_x_inches),
            Inches(callout_y_inches),
            Inches(120 / 144),  # 120px at 144 DPI
            Inches(40 / 144)    # 40px at 144 DPI
        )

        # Style callout
        callout_box.fill.solid()
        callout_box.fill.fore_color.rgb = RGBColor(255, 255, 224)  # Light yellow
        callout_box.line.color.rgb = RGBColor(20, 123, 88)  # Primary Green
        callout_box.line.width = Pt(1)

        # Add text
        callout_frame = callout_box.text_frame
        callout_frame.text = annotation.get('text', '')
        callout_frame.word_wrap = True
        callout_frame.margin_left = px_to_inches(5)
        callout_frame.margin_right = px_to_inches(5)
        callout_frame.margin_top = px_to_inches(3)
        callout_frame.margin_bottom = px_to_inches(3)

        p = callout_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(74, 74, 74)

    def _add_key_takeaway_box(self, slide, text: str, x, y, width=None, height=None):
        """
        IMPROVEMENT 5: Add a "Key Takeaway" callout box for storytelling.

        This is a bordered box that summarizes the most critical insight on the slide.
        Styled with:
        - Light grey background (#F5F5F5)
        - 2pt Primary Green border (#147B58)
        - Contains single, impactful sentence
        - Positioned in white space area (often within chart area)

        Args:
            slide: PowerPoint slide object
            text: The key takeaway message (single impactful sentence)
            x, y: Position in inches
            width, height: Size in inches (defaults to auto-fit based on text)
        """
        # Default dimensions if not specified
        if width is None:
            width = px_to_inches(400)  # ~400px wide by default
        if height is None:
            height = px_to_inches(60)  # ~60px tall by default

        # Create callout box
        takeaway_box = slide.shapes.add_textbox(x, y, width, height)

        # IMPROVEMENT 5: Styling - Light grey background
        takeaway_box.fill.solid()
        takeaway_box.fill.fore_color.rgb = RGBColor(245, 245, 245)  # #F5F5F5

        # IMPROVEMENT 5: Styling - 2pt Primary Green border
        takeaway_box.line.color.rgb = RGBColor(*self.colors['primary_green_rgb'])
        takeaway_box.line.width = Pt(2)

        # Add text with professional formatting
        text_frame = takeaway_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_left = px_to_inches(12)
        text_frame.margin_right = px_to_inches(12)
        text_frame.margin_top = px_to_inches(8)
        text_frame.margin_bottom = px_to_inches(8)

        # Format paragraph
        p = text_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = 'Arial'
        p.font.size = Pt(12)  # Readable size
        p.font.bold = True  # Bold for impact
        p.font.color.rgb = RGBColor(*self.colors['primary_green_rgb'])  # Primary Green text

        return takeaway_box

    def _validate_element_bounds(self, x_px: float, y_px: float, width_px: float, height_px: float) -> bool:
        """
        Validate that element is within Safe Zone boundaries.

        Args:
            x_px: X position in pixels
            y_px: Y position in pixels
            width_px: Width in pixels
            height_px: Height in pixels

        Returns:
            True if element is within Safe Zone, False otherwise
        """
        return validate_bounds(x_px, y_px, width_px, height_px)

    def get_config(self) -> Dict[str, Any]:
        """Get Main v2.0 configuration."""
        return self.config

    def get_grid_info(self) -> Dict[str, Any]:
        """Get grid system information."""
        return {
            'canvas': self.grid['canvas'],
            'safe_zone': self.grid['safe_zone'],
            'regions': self.grid['regions']
        }
